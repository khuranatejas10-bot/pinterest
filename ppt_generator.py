import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_pinterest_ppt(keyword, image_paths, output_path):
    """
    Generates a PowerPoint presentation using the collected Pinterest images.
    
    :param keyword: The query keyword searched
    :param image_paths: List of local paths to the downloaded images (max 10)
    :param output_path: Where to save the generated .pptx file
    """
    prs = Presentation()
    # Use standard 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide (Dark Background)
    blank_slide_layout = prs.slide_layouts[6]  # Completely blank slide layout
    slide1 = prs.slides.add_slide(blank_slide_layout)
    
    # Set dark background (#0F0F0F)
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 15, 15)
    
    # Large bold white centered keyword text
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = keyword.upper()
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    
    # Pinterest red accent line (#E60023)
    # add_shape parameters: type, left, top, width, height
    accent_bar = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.666), Inches(3.7), Inches(4), Inches(0.08)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor(230, 0, 35) # #E60023
    accent_bar.line.color.rgb = RGBColor(230, 0, 35)
    
    # Smaller label "Top 10 on Pinterest"
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.333), Inches(1))
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Top 10 on Pinterest"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(180, 180, 180) # light gray
    p2.font.name = 'Arial'
    p2.alignment = PP_ALIGN.CENTER
    
    # Slides 2 through 6 (Side-by-side images, 2 per slide)
    num_images = len(image_paths)
    
    # Calculate how many slides we actually need based on image count
    # 10 images = 5 slides. If we have fewer, we still want to make as many as possible.
    num_data_slides = (num_images + 1) // 2
    if num_data_slides == 0:
        # Save empty ppt if no images
        prs.save(output_path)
        return 0, 0
    
    # Dimensions for side-by-side images
    # Slide width = 13.333, height = 7.5
    # Margins: Left/Right/Top/Bottom = 0.5 inches. Padding between images = 0.5 inches
    img_width = Inches(5.9)
    img_height = Inches(6.1)
    img_top = Inches(0.7)
    
    left_positions = [Inches(0.5), Inches(6.9)]
    
    for i in range(num_images):
        slide_idx = i // 2
        col_idx = i % 2
        
        # Check if we need to create a new slide
        if col_idx == 0:
            current_slide = prs.slides.add_slide(blank_slide_layout)
            # Add a subtle light/dark background to slide 2-6 (let's make it a clean dark background too, or a very dark gray)
            bg = current_slide.background
            bg_fill = bg.fill
            bg_fill.solid()
            bg_fill.fore_color.rgb = RGBColor(24, 24, 24) # #181818
            
        img_path = image_paths[i]
        img_left = left_positions[col_idx]
        
        # Add the image shape
        try:
            current_slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)
        except Exception as e:
            # If python-pptx fails to parse the image, let's put a placeholder text or skip it
            print(f"Error adding picture {img_path}: {e}")
            # Add a placeholder box
            placeholder = current_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                img_left, img_top, img_width, img_height
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(40, 40, 40)
            tf_ph = placeholder.text_frame
            tf_ph.word_wrap = True
            p_ph = tf_ph.paragraphs[0]
            p_ph.text = "[Image Unavailable]"
            p_ph.font.size = Pt(18)
            p_ph.font.color.rgb = RGBColor(150, 150, 150)
            p_ph.alignment = PP_ALIGN.CENTER
            
        # Add rank badge in the bottom-right corner of the image
        # Size: 0.6 x 0.6 inches
        badge_size = Inches(0.6)
        badge_left = img_left + img_width - badge_size - Inches(0.15)
        badge_top = img_top + img_height - badge_size - Inches(0.15)
        
        badge = current_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            badge_left, badge_top, badge_size, badge_size
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(230, 0, 35) # Pinterest red
        badge.line.color.rgb = RGBColor(230, 0, 35)
        
        # Badge text (1-based rank number)
        tf_badge = badge.text_frame
        p_badge = tf_badge.paragraphs[0]
        p_badge.text = str(i + 1)
        p_badge.font.size = Pt(16)
        p_badge.font.bold = True
        p_badge.font.color.rgb = RGBColor(255, 255, 255)
        p_badge.alignment = PP_ALIGN.CENTER
        
    prs.save(output_path)
    return num_images, num_data_slides + 1 # Return number of images and total slides generated
